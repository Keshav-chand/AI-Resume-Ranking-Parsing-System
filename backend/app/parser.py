import io
import re
from pathlib import Path

import pdfplumber
from docx import Document
from pptx import Presentation


def extract_text_from_pdf(file_bytes: bytes) -> str:
    text = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
    return "\n".join(text)


def extract_text_from_docx(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    paragraphs.append(cell.text.strip())
    return "\n".join(paragraphs)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
    return "\n".join(text_parts)


def parse_file(filename: str, file_bytes: bytes) -> str:
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext in (".doc", ".docx"):
        return extract_text_from_docx(file_bytes)
    elif ext in (".ppt", ".pptx"):
        return extract_text_from_pptx(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def clean_text(text: str) -> str:
    text = re.sub(r"\r\n|\r", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()